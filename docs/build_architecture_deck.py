"""Generate the Ed-Fi Text-to-SQL client-facing architecture deck.

Run with:
    uvx --from python-pptx python docs/build_architecture_deck.py

Produces docs/architecture_deck.pptx.

Audience: client / business stakeholder. Frame the work in outcomes,
not file paths. Status is grounded in the actual repo audit (Apr 30
2026): 13 core components built and tested; agent loop and frontend
shipped; remaining 6 working weeks target Leiden sub-clustering,
multi-provider matrix tests, agentic polish, observability,
performance, documentation, and a recorded demo for the v0.9
release on Fri June 13.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Palette ────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x14, 0x1A)
PANEL    = RGBColor(0x1A, 0x21, 0x2A)
BORDER   = RGBColor(0x2A, 0x33, 0x3F)
TEXT     = RGBColor(0xE6, 0xEA, 0xF0)
MUTED    = RGBColor(0x8A, 0x95, 0xA5)
ACCENT   = RGBColor(0x6E, 0xC1, 0xFF)
GREEN    = RGBColor(0x4A, 0xD0, 0x9C)
AMBER    = RGBColor(0xF2, 0xB8, 0x57)
RED      = RGBColor(0xE5, 0x6B, 0x6B)


# ── Helpers ────────────────────────────────────────────────────────────────


def _set_text(tf, text, *, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Helvetica Neue"


def _add_text(slide, text, *, left, top, width, height,
              size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    tf.word_wrap = True
    _set_text(tf, text, size=size, bold=bold, color=color, align=align)
    return box


def _bullets(slide, items, *, left, top, width, height, size=14, color=TEXT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.1)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = "• " + item if not item.startswith(("• ", "  ")) else item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Helvetica Neue"


def _panel(slide, *, left, top, width, height, fill=PANEL, border=BORDER):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def _slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _add_title(slide, title, subtitle=None):
    _add_text(slide, title, left=0.5, top=0.35, width=12.3, height=0.7,
              size=30, bold=True, color=ACCENT)
    if subtitle:
        _add_text(slide, subtitle, left=0.5, top=1.05, width=12.3, height=0.45,
                  size=14, color=MUTED)


def _footer(slide, left_text, right_text):
    _add_text(slide, left_text, left=0.5, top=7.05, width=8, height=0.3,
              size=10, color=MUTED)
    _add_text(slide, right_text, left=8.5, top=7.05, width=4.3, height=0.3,
              size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _status_chip(slide, *, left, top, label, color, width=1.25):
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top), Inches(width), Inches(0.32))
    chip.fill.solid(); chip.fill.fore_color.rgb = color
    chip.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.32))
    tf = tb.text_frame; tf.margin_left = tf.margin_right = Inches(0); tf.margin_top = Inches(0.04)
    _set_text(tf, label, size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)


def _section_divider(slide, kicker, headline, subhead, color=ACCENT):
    _add_text(slide, kicker, left=0.8, top=2.5, width=11.7, height=0.5,
              size=18, color=MUTED)
    _add_text(slide, headline, left=0.8, top=3.05, width=11.7, height=1.0,
              size=42, bold=True, color=color)
    _add_text(slide, subhead, left=0.8, top=4.05, width=11.7, height=0.5,
              size=18, color=TEXT)


# ── Deck construction ─────────────────────────────────────────────────────


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Slide 1: Title ────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_text(s, "Ed-Fi Text-to-SQL", left=0.8, top=2.1, width=11.7, height=1.0,
              size=46, bold=True, color=ACCENT)
    _add_text(s, "Ask questions of your school data in plain English",
              left=0.8, top=3.15, width=11.7, height=0.6,
              size=22, color=TEXT)
    _add_text(s, "Architecture review · platform built · 6 working weeks to v0.9 release",
              left=0.8, top=3.85, width=11.7, height=0.5,
              size=16, color=MUTED)
    _add_text(s, "Apr 30, 2026  ·  client review  ·  v0.9 release Fri June 13",
              left=0.8, top=6.6, width=11.7, height=0.4,
              size=13, color=MUTED)
    _notes(s, "Set the room: this is an architecture-and-status review. The platform's foundation "
              "is built — knowledge layer, NL→SQL pipeline, agent loop, API, frontend, and ops "
              "tooling are all functionally complete and tested. The next six working weeks "
              "(Apr 30 through Fri June 12, weekends + May 1 holiday excluded) close the loop on "
              "agentic polish, multi-provider hardening, observability, performance, recorded "
              "demo, and documentation. v0.9 release is tagged Fri June 13. The honest reason "
              "the timeline runs to mid-June: the punch list is real work — provider matrix tests "
              "alone span 5 LLMs × 4 embeddings × 3 dialects × 4 consumers — and we've baked in "
              "buffer for the inevitable surprises rather than booking a heroics-only plan.")

    # ── Slide 2: The problem ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_title(s, "The problem we're solving",
                  "An analyst asks a question. The platform answers with data.")

    _panel(s, left=0.5, top=1.6, width=12.3, height=2.1, fill=PANEL)
    _add_text(s, "Today, in any Ed-Fi-shaped database",
              left=0.7, top=1.7, width=11.9, height=0.4, size=14, bold=True, color=AMBER)
    _bullets(s, [
        "1,048 tables · ~10,000 columns · 1,663 foreign keys (Ed-Fi DS 6.1.0)",
        "An analyst with a question must (a) know which tables to join, (b) write the SQL, (c) execute it, (d) read it",
        "End-to-end this is hours of work for a senior engineer — and impossible for a non-engineer",
    ], left=0.7, top=2.1, width=11.9, height=1.6, size=13)

    _panel(s, left=0.5, top=3.9, width=12.3, height=2.9, fill=PANEL)
    _add_text(s, "What we're delivering by Fri June 13",
              left=0.7, top=4.0, width=11.9, height=0.4, size=14, bold=True, color=GREEN)
    _bullets(s, [
        "Type a question in English  →  receive correct SQL, the rows it returns, a chart, and a written summary",
        "Have a multi-turn conversation with the platform — it remembers what you asked before",
        "Works on any Ed-Fi database (Postgres / MSSQL / SQLite) without changing the database itself",
        "Cites the tables it used so an engineer can verify the answer in seconds",
        "Learns from approved queries — the more it's used, the better it gets",
    ], left=0.7, top=4.4, width=11.9, height=2.4, size=13)

    _footer(s, "Why this project · Section 1", "")
    _notes(s, "Anchor in the user pain. Most Ed-Fi work today goes through a small handful of "
              "people who can write SQL against the model. We're widening that bottleneck — not "
              "by training more SQL writers but by automating the translation from English. "
              "Mention the conversation bullet — that's the agent-loop work that closes in this "
              "release.")

    # ── Slide 3: Status snapshot ──────────────────────────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_title(s, "Status: foundation built, polish ahead",
                  "13 core components green and tested · 4 workstreams remaining for v0.9 on Fri June 13")

    rows = [
        ("Knowledge layer",
         "Done", GREEN,
         "Domain classification (829/829 tables), FK graph (1,663 edges), APSP, Steiner solver, table catalog, hybrid retrieval, 4-tier entity resolver"),
        ("NL → SQL pipeline",
         "Done", GREEN,
         "Routed orchestrator: classify → retrieve → resolve → generate → validate → repair → execute → chart + summary"),
        ("Agentic layer",
         "Done", GREEN,
         "8-tool agent loop, multi-turn conversations with cross-dialect history, server-sent streaming"),
        ("API + Web UI",
         "Done", GREEN,
         "18 endpoints (REST + WebSocket); 7 pages — Query, Chat, Tables, Domains, Gold, Settings, browser; provider-aware refresh"),
        ("Operations",
         "Done", GREEN,
         "Multi-dialect metadata DB, runtime Settings UI, eval harness with JSON+Markdown reports, ships zero-infra demo SQLite"),
        ("Auto sub-clustering (Leiden)",
         "In flight", AMBER,
         "Refines oversize Ed-Fi domains into 8–20-table clusters for sharper routing; static taxonomy works today"),
        ("Multi-provider matrix tests",
         "In flight", AMBER,
         "5 LLM providers + 4 embedding providers wired; cross-product regression matrix lands by June 13"),
        ("Agentic polish",
         "In flight", AMBER,
         "Tool reliability, streaming UX refinements, more agent tools (chart export, eval drilldown)"),
        ("Observability (OTel + Prometheus)",
         "Coming", RED,
         "Per-stage spans + counters; logging is in place today, traces and metrics are next"),
    ]
    _panel(s, left=0.5, top=1.6, width=12.3, height=5.4, fill=PANEL)
    y = 1.72
    for label, status, color, what in rows:
        _add_text(s, label, left=0.7, top=y, width=4.3, height=0.32,
                  size=12, bold=True, color=TEXT)
        _status_chip(s, left=5.1, top=y + 0.02, label=status, color=color, width=1.15)
        _add_text(s, what, left=6.45, top=y + 0.02, width=6.3, height=0.6,
                  size=10, color=MUTED)
        y += 0.58

    _footer(s, "Status snapshot", "v0.9 release · Fri June 13, 2026")
    _notes(s, "This is the headline. Five rows green = foundation done; three amber = active work "
              "the team finishes by June 13; one red = stretch item that may slip into v0.10. "
              "The honest 6-week timeline (vs the original May 8) reflects the real punch list: "
              "provider matrix tests across the full 5 × 4 × 3 × 4 grid, Leiden tuning that "
              "needs iteration, observability work that touches every stage, and a recorded demo "
              "we want to be able to ship without disclaimers.")

    # ── Slide 4: System map ───────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_title(s, "The big picture",
                  "Five layers carry every question from English to answer")

    layers = [
        ("Question",       ["Analyst types: 'How many Hispanic students enrolled in Grade 9 last year?'"], MUTED),
        ("Knowledge",      ["Classification routes the question to the right Ed-Fi domains",
                            "Hybrid retrieval narrows 1,048 tables to a small relevant set",
                            "FK graph + Steiner returns the cheapest JOIN tree connecting them",
                            "Entity resolver maps real-world terms (Hispanic, Grade 9) to codes"], GREEN),
        ("Pipeline",       ["LLM writes SQL grounded in real columns; validator parses + plans + dry-runs",
                            "Repair loop (3 attempts) auto-fixes the common failure modes"], GREEN),
        ("Agentic",        ["Agent loop calls 8 tools across multi-turn conversations",
                            "Server-sent streaming gives token-level feedback in the browser"], GREEN),
        ("Answer",         ["SQL · rows · auto-picked chart · plain-English summary · cited tables"], GREEN),
    ]
    y = 1.4
    for label, lines, color in layers:
        h = 0.55 + 0.32 * len(lines)
        _panel(s, left=0.5, top=y, width=12.3, height=h, fill=PANEL)
        _add_text(s, label, left=0.7, top=y + 0.1, width=2.6, height=0.4,
                  size=13, bold=True, color=color)
        for j, ln in enumerate(lines):
            _add_text(s, "• " + ln, left=3.4, top=y + 0.1 + j * 0.30, width=9.2, height=0.32,
                      size=11, color=TEXT)
        y += h + 0.08

    _add_text(s, "Every layer is built and tested today; the work to Fri June 13 is polish, hardening, and recorded demo",
              left=0.5, top=7.05, width=12.3, height=0.3,
              size=11, color=MUTED, align=PP_ALIGN.CENTER)
    _notes(s, "Walk top-to-bottom. Each row corresponds to a tested module set in the repo. "
              "Spend the most time on Pipeline + Agentic — those are the parts the demo audience "
              "actually sees in motion.")

    # ── Slide 5: By the numbers ───────────────────────────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_title(s, "By the numbers",
                  "What the platform actually contains today")

    cols = [
        ("Knowledge", GREEN, [
            ("Tables modeled",            "1,048"),
            ("Foreign keys parsed",        "1,663"),
            ("Domains (Ed-Fi taxonomy)",      "35"),
            ("Catalog entries",             "829"),
            ("APSP matrix",          "829 × 829"),
            ("Steiner p99 latency",         "< 1 ms"),
        ]),
        ("Pipeline", GREEN, [
            ("Pipeline stages end-to-end",     "9"),
            ("Repair attempts before fail",    "3"),
            ("Entity-resolver tiers",          "4"),
            ("Hybrid retrieval blend",  "0.6 / 0.4"),
            ("Gold-store tests passing",     "7/7"),
            ("Pipeline tests passing",       "4/4"),
        ]),
        ("Surface", GREEN, [
            ("API endpoints",                 "18"),
            ("WebSocket streaming routes",     "2"),
            ("Frontend pages",                 "7"),
            ("Agent tools wired",              "8"),
            ("LLM providers integrated",       "5"),
            ("Database dialects supported",    "3"),
        ]),
    ]
    x = 0.5
    for header, color, rows in cols:
        _panel(s, left=x, top=1.6, width=4.1, height=5.3, fill=PANEL)
        _add_text(s, header, left=x + 0.2, top=1.7, width=3.7, height=0.4,
                  size=15, bold=True, color=color)
        y = 2.2
        for k, v in rows:
            _add_text(s, k, left=x + 0.2, top=y, width=2.5, height=0.32,
                      size=11, color=MUTED)
            _add_text(s, v, left=x + 2.6, top=y, width=1.4, height=0.32,
                      size=12, bold=True, color=TEXT, align=PP_ALIGN.RIGHT)
            y += 0.55
        x += 4.25

    _footer(s, "Numbers are from the latest test run + Ed-Fi DS 6.1.0 build", "")
    _notes(s, "These are the real numbers from the repo audit. Use them to anchor any 'is this "
              "actually working?' question — every figure has a corresponding test or artifact.")

    # ── Slide 6: Section divider — Knowledge ──────────────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _section_divider(s,
        "Section 2 of 5",
        "The Knowledge Layer",
        "How the platform turns 1,048 tables into a navigable, queryable index",
        color=GREEN)
    _notes(s, "Foundation. Without this layer the LLM would be guessing across 1,048 tables.")

    # ── Slide 7: Domain classification + table catalog ────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_title(s, "Cataloging every Ed-Fi table",
                  "Four-stage classifier · 829/829 tables placed · descriptions captured at build time")

    _panel(s, left=0.5, top=1.6, width=6.0, height=5.3, fill=PANEL)
    _add_text(s, "Four-stage classifier", left=0.7, top=1.7, width=5.6, height=0.4,
              size=15, bold=True, color=ACCENT)
    rows = [
        ("Stage 1 · Direct ApiModel domains",  "750 tables"),
        ("Stage 2 · Aggregate inheritance",      "75 tables"),
        ("Stage 3 · Descriptor FK voting",        "3 tables"),
        ("Stage 4 · LLM fallback",         "rare residual"),
        ("Override layer (operator)",     "human in the loop"),
        ("Output artifact",      "table_classification.json"),
    ]
    y = 2.2
    for k, v in rows:
        _add_text(s, k, left=0.7, top=y, width=3.5, height=0.32, size=12, color=MUTED)
        _add_text(s, v, left=4.2, top=y, width=2.0, height=0.32,
                  size=12, bold=True, color=TEXT, align=PP_ALIGN.RIGHT)
        y += 0.58

    _panel(s, left=6.85, top=1.6, width=6.0, height=5.3, fill=PANEL)
    _add_text(s, "Per-table catalog entry",
              left=7.05, top=1.7, width=5.6, height=0.4, size=15, bold=True, color=AMBER)
    _bullets(s, [
        "Authoritative description (Ed-Fi prose, no LLM hallucination)",
        "Domain tags + identifying columns + linked descriptors",
        "Connected tables (FK neighbors) for join intuition",
        "Sample column values for entity resolution",
        "Proven-query count — biases retrieval to tables we've answered before",
        "Reflects unknown tables on the live DB so non-Ed-Fi columns join cleanly",
    ], left=7.05, top=2.15, width=5.6, height=4.5, size=12)

    _footer(s, "Knowledge · Classification & Catalog", "test_component2 + test_component4 green")
    _notes(s, "Stage 1 handles 90% of tables deterministically. The LLM tier is the safety net "
              "for residual stragglers. Reflection of unknown DB tables (P7) handles the case "
              "where the operator's database has 36 extra tables that aren't in the Ed-Fi spec.")

    # ── Slide 8: Ingestion flow when target DB diverges from standard Ed-Fi ─
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_title(s, "Ingesting a non-standard target database",
                  "Five-step flow when the operator's DB has tables Ed-Fi never specified · "
                  "e.g. 1,084 in DB vs 1,048 in spec")

    steps = [
        ("1",
         "Build standard Ed-Fi knowledge (one-time per DS version)",
         "Fetch ApiModel.json + 0030-ForeignKeys.sql from Ed-Fi GitHub raw · cache under data/edfi/ · "
         "4-stage classifier on the full ApiModel produces 829 catalog entries and 1,663 FK edges",
         GREEN),
        ("2",
         "Connect & activate the operator's database",
         "Operator opens Settings → Database connector form · picks dialect (Postgres / MSSQL / SQLite) "
         "and enters host + credentials (passwords stored in gitignored runtime_secrets.json) · "
         "Make active creates a per-provider artifact directory",
         ACCENT),
        ("3",
         "Reflect tables Ed-Fi doesn't know about",
         "SQLAlchemy Inspector walks the live DB · for each table NOT in ApiModel: reflect columns + types + PK + FKs "
         "(composite supported) · sample N rows · tag domains=['Other'], is_extension=True · LLM gap-fill writes a description",
         ACCENT),
        ("4",
         "Build per-provider artifacts (no cross-pollution)",
         "Catalog merges 829 Ed-Fi entries + N reflected entries · graph merges 1,663 Ed-Fi edges + reflected FK edges so "
         "cross-source joins emerge naturally · embeddings re-indexed · everything stored under "
         "data/artifacts/per_provider/<provider_name>/",
         ACCENT),
        ("5",
         "Query routes through the unified catalog",
         "Routing matches Ed-Fi domains OR the 'Other' bucket · hybrid retrieval scores Ed-Fi + reflected tables together · "
         "Steiner picks join paths across both sources · LLM grounds SQL in real columns from whichever tables matter",
         GREEN),
    ]
    _panel(s, left=0.5, top=1.6, width=12.3, height=5.4, fill=PANEL)
    y = 1.75
    for num, title, body, color in steps:
        # Numbered badge
        chip = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(0.7), Inches(y), Inches(0.55), Inches(0.55))
        chip.fill.solid(); chip.fill.fore_color.rgb = color
        chip.line.fill.background()
        tb = s.shapes.add_textbox(Inches(0.7), Inches(y), Inches(0.55), Inches(0.55))
        tf = tb.text_frame; tf.margin_left = tf.margin_right = Inches(0); tf.margin_top = Inches(0.1)
        _set_text(tf, num, size=18, bold=True, color=BG, align=PP_ALIGN.CENTER)
        # Title + body
        _add_text(s, title, left=1.4, top=y - 0.02, width=11.4, height=0.32,
                  size=13, bold=True, color=TEXT)
        _add_text(s, body, left=1.4, top=y + 0.32, width=11.4, height=0.65,
                  size=10, color=MUTED)
        y += 1.04

    _footer(s, "Knowledge · Non-standard ingest flow",
            "test_catalog_reflect_unknown.py · 8 cases green")
    _notes(s, "This is the answer to 'what if my database has tables that aren't in the Ed-Fi spec?' "
              "Step 3 (reflection) is the key insight — we don't require the operator to write a "
              "schema-mapping config. SQLAlchemy's Inspector walks the live DB, anything we don't "
              "recognize gets reflected with composite FKs and sample rows, then merges into the "
              "same per-provider catalog the Ed-Fi tables live in. Per-provider artifact isolation "
              "(N1-N5) means switching from one operator DB to another flips a different set of "
              "files in — no stale state from the previous DB leaks into the next one.")

    # ── Slide 9 (was 8): FK graph + Steiner ───────────────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_title(s, "Foreign-Key Graph + JOIN solver",
                  "1,663 edges · 829×829 APSP precomputed · KMB Steiner approximation")

    _panel(s, left=0.5, top=1.6, width=6.0, height=5.3, fill=PANEL)
    _add_text(s, "What the graph contains", left=0.7, top=1.7, width=5.6, height=0.4,
              size=15, bold=True, color=ACCENT)
    rows = [
        ("Nodes (tables)",                      "829"),
        ("Edges (FK relationships)",          "1,663"),
        ("APSP storage on disk",               "<10 MB"),
        ("Load time at startup",              "<100 ms"),
        ("Same-aggregate join weight",            "1.0"),
        ("Cross-aggregate join weight",           "2.0"),
        ("Cross-domain join weight",              "4.0"),
        ("Composite-FK bonus multiplier",         "0.9"),
    ]
    y = 2.15
    for k, v in rows:
        _add_text(s, k, left=0.7, top=y, width=3.7, height=0.32, size=12, color=MUTED)
        _add_text(s, v, left=4.4, top=y, width=1.8, height=0.32,
                  size=12, bold=True, color=TEXT, align=PP_ALIGN.RIGHT)
        y += 0.50

    _panel(s, left=6.85, top=1.6, width=6.0, height=5.3, fill=PANEL)
    _add_text(s, "Steiner solver — three paths",
              left=7.05, top=1.7, width=5.6, height=0.4, size=15, bold=True, color=AMBER)
    rows = [
        ("k = 2 tables",
         "Bidirectional Dijkstra; returns top-3 alternatives so the model picks the most readable",
         "<5 ms"),
        ("k = 3–8 tables",
         "KMB Steiner-tree approximation — provably within 2× of optimal",
         "<50 ms"),
        ("k > 8 (edge case)",
         "Greedy fallback; rare on real Ed-Fi questions",
         "<200 ms"),
    ]
    y = 2.15
    for label, what, perf in rows:
        _add_text(s, label, left=7.05, top=y, width=5.6, height=0.32,
                  size=12, bold=True, color=TEXT)
        _add_text(s, what, left=7.05, top=y + 0.30, width=5.6, height=0.7, size=11, color=MUTED)
        _add_text(s, perf, left=11.0, top=y, width=1.7, height=0.32,
                  size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
        y += 1.30

    _footer(s, "Knowledge · FK Graph & Steiner",
            "test_component3 + test_component3_benchmark green · p99 < 1 ms")
    _notes(s, "Two ideas to leave the room with: graph is small (10 MB) and fast (sub-millisecond "
              "responses), and weights bias joins to stay inside the natural Ed-Fi entity "
              "boundaries — which is what a senior analyst would do by hand.")

    # ── Slide 9: Embedding + retrieval + entity resolution ────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_title(s, "Finding the right tables, resolving real-world terms",
                  "Hybrid retrieval shortlists tables · 4-tier resolver maps user words to codes")

    _panel(s, left=0.5, top=1.6, width=6.0, height=5.3, fill=PANEL)
    _add_text(s, "Hybrid retrieval (semantic + keyword)",
              left=0.7, top=1.7, width=5.6, height=0.4, size=15, bold=True, color=ACCENT)
    _bullets(s, [
        "Per-table semantic blob: description, key columns, neighbors, linked descriptors",
        "Embedded once at build time; queried millions of times after",
        "Domain pre-filter narrows the candidate pool before scoring",
        "Cosine similarity (0.6) + BM25 keyword score (0.4) — best of both",
        "Returns the small set (typically 3–8 tables) that best matches the question",
    ], left=0.7, top=2.15, width=5.6, height=4.5, size=12)

    _panel(s, left=6.85, top=1.6, width=6.0, height=5.3, fill=PANEL)
    _add_text(s, "Four-tier entity resolver",
              left=7.05, top=1.7, width=5.6, height=0.4, size=15, bold=True, color=AMBER)
    rows = [
        ("Tier 1 · Exact lookup",  "Hash hit on the value index", "<1 ms"),
        ("Tier 2 · Fuzzy match",    "Typos: 'Hispanc' → 'Hispanic'", "5 ms"),
        ("Tier 3 · Semantic ANN",  "Match meaning when spelling differs", "20 ms"),
        ("Tier 4 · LLM disambiguate","Resolve when one term means two things", "200 ms"),
    ]
    y = 2.2
    for tier, what, perf in rows:
        _add_text(s, tier, left=7.05, top=y, width=3.6, height=0.32, size=12, bold=True, color=TEXT)
        _add_text(s, what, left=7.05, top=y + 0.30, width=4.4, height=0.32, size=11, color=MUTED)
        _add_text(s, perf, left=11.4, top=y, width=1.3, height=0.32,
                  size=11, color=GREEN, align=PP_ALIGN.RIGHT)
        y += 0.85

    _footer(s, "Knowledge · Retrieval & Entity Resolution",
            "test_component5 + test_component6 green · Tier 1+2 handle ~95% of cases")
    _notes(s, "This is where 'plain English' becomes 'database-correct'. Without it, "
              "questions like 'last year', 'Grade 9', or 'free-and-reduced lunch' would never "
              "land on the right rows.")

    # ── Slide 10: Section divider — Pipeline ──────────────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _section_divider(s,
        "Section 3 of 5",
        "The NL → SQL Pipeline",
        "Nine stages, end-to-end · validation + 3-attempt repair · charts and descriptions",
        color=GREEN)
    _notes(s, "This is the orchestrator that knits the knowledge layer into an answer.")

    # ── Slide 11: Pipeline orchestrator ───────────────────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_title(s, "The nine pipeline stages",
                  "answer(nl) → SQL + rows + chart + summary · all wired and tested today")

    stages = [
        ("1. Classify",       "Route question → primary + secondary + tertiary domains"),
        ("2. Retrieve",       "Hybrid search inside routed domains → 3–8 candidate tables"),
        ("3. Resolve",        "Map user terms to descriptor codes (4-tier funnel)"),
        ("4. Recall gold",    "Top-3 proven NL/SQL pairs by AST + NL similarity"),
        ("5. Generate",       "LLM writes SQL grounded in real columns + JOINs"),
        ("6. Validate",       "sqlglot parse + EXPLAIN + LIMIT 0 dry-run"),
        ("7. Repair",         "3-attempt loop fixes typos / missing joins / wrong columns"),
        ("8. Execute",        "Run against the active target DB engine"),
        ("9. Visualize + describe", "Vega-Lite chart + plain-English summary in parallel"),
    ]
    _panel(s, left=0.5, top=1.6, width=12.3, height=5.4, fill=PANEL)
    y = 1.78
    for label, what in stages:
        _add_text(s, label, left=0.7, top=y, width=3.7, height=0.32,
                  size=13, bold=True, color=ACCENT)
        _add_text(s, what, left=4.5, top=y, width=8.2, height=0.32,
                  size=12, color=TEXT)
        y += 0.56

    _footer(s, "Pipeline · 9 stages · test_component8 green",
            "Postgres / MSSQL / SQLite all pass")
    _notes(s, "This slide is the heart of the demo. Walk through each stage with the audience "
              "imagining their own question — by the end they should feel like they've watched "
              "the platform think.")

    # ── Slide 12: Validation, repair, charts, descriptions ────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_title(s, "Self-correcting and self-explaining",
                  "Bad SQL never reaches production · charts and summaries land in parallel")

    _panel(s, left=0.5, top=1.6, width=6.0, height=5.3, fill=PANEL)
    _add_text(s, "Validation + repair loop",
              left=0.7, top=1.7, width=5.6, height=0.4, size=15, bold=True, color=ACCENT)
    _bullets(s, [
        "Parse-check — catches syntax errors (sqlglot)",
        "Plan-check — catches semantic errors (EXPLAIN)",
        "Dry-execute — catches runtime errors (LIMIT 0)",
        "Up to three repair attempts with fresh context",
        "Live-fix verified end-to-end (test_repair_fixes_real_explain_error)",
        "Short-circuits on clean first SQL — no overhead in the common case",
    ], left=0.7, top=2.15, width=5.6, height=4.5, size=12)

    _panel(s, left=6.85, top=1.6, width=6.0, height=5.3, fill=PANEL)
    _add_text(s, "Charts + descriptions (parallel)",
              left=7.05, top=1.7, width=5.6, height=0.4, size=15, bold=True, color=AMBER)
    _bullets(s, [
        "Result-shape inference: row count, column types, aggregates, temporal axes",
        "Vega-Lite spec generated via LLM with strict JSON schema",
        "Plain-English summary generated in parallel via ThreadPoolExecutor",
        "Browser receives chart spec + summary the moment they're ready",
        "Single-row, multi-row, time-series, and aggregate cases all handled",
    ], left=7.05, top=2.15, width=5.6, height=4.5, size=12)

    _footer(s, "Pipeline · Validation & Output",
            "test_component9 + test_component10 green")
    _notes(s, "Repair is the magic that lets the platform be honest with the user. If the LLM "
              "writes a wrong column, EXPLAIN catches it; the system asks the LLM to fix it with "
              "the actual error message; we get the right answer most of the time. This is what "
              "makes the platform safe for non-engineers.")

    # ── Slide 13: Section divider — Agentic ───────────────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _section_divider(s,
        "Section 4 of 5",
        "The Agentic Layer",
        "8 tools · multi-turn conversations · server-sent streaming",
        color=GREEN)
    _notes(s, "The agent loop is what makes this feel like a partner, not a one-shot translator.")

    # ── Slide 14: Agent loop ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_title(s, "Eight tools the agent can call",
                  "Each tool is the same primitive an engineer would reach for — exposed to the LLM")

    tools = [
        ("CLASSIFY_DOMAINS",     "Find the Ed-Fi domains relevant to a question"),
        ("SEARCH_TABLES",         "Hybrid search inside the catalog · top-K tables"),
        ("INSPECT_TABLE",         "Full schema + sample rows for one table"),
        ("RESOLVE_ENTITY",        "Map a user term ('Hispanic') to (table, column, code)"),
        ("FIND_JOIN_PATH",        "Steiner solver between two or more tables"),
        ("FIND_SIMILAR_QUERIES", "Top-3 proven NL/SQL pairs from the gold store"),
        ("RUN_SQL",               "Execute SQL against the active target DB · returns rows"),
        ("FINAL_ANSWER",          "Terminate the turn with a structured answer"),
    ]
    _panel(s, left=0.5, top=1.6, width=12.3, height=5.4, fill=PANEL)
    y = 1.8
    for tool, what in tools:
        _add_text(s, tool, left=0.7, top=y, width=4.0, height=0.32,
                  size=12, bold=True, color=ACCENT)
        _add_text(s, what, left=4.8, top=y, width=8.0, height=0.32,
                  size=12, color=TEXT)
        y += 0.62

    _footer(s, "Agentic · 8 tools wired · max 20 steps per turn",
            "test_agent_loop · streaming events verified")
    _notes(s, "The agent never invents tables or values — it must call SEARCH_TABLES and "
              "RESOLVE_ENTITY first. RUN_SQL is the only side-effecting tool, and it executes "
              "validated SQL only.")

    # ── Slide 15: Conversations + streaming + providers ───────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_title(s, "Conversations, streaming, and provider portability",
                  "The agent works the same on Anthropic, Bedrock, Azure, OpenAI, OpenRouter")

    _panel(s, left=0.5, top=1.6, width=6.0, height=5.3, fill=PANEL)
    _add_text(s, "Multi-turn conversations",
              left=0.7, top=1.7, width=5.6, height=0.4, size=15, bold=True, color=ACCENT)
    _bullets(s, [
        "SQLAlchemy-backed conversation store — Postgres, MSSQL, SQLite",
        "Each conversation tagged with the dialect it was authored in",
        "Cross-dialect history visible (badge in the chat UI)",
        "Server-sent events stream tokens, tool calls, and final answers",
        "Frontend reducer renders updates without page reloads",
    ], left=0.7, top=2.15, width=5.6, height=4.5, size=12)

    _panel(s, left=6.85, top=1.6, width=6.0, height=5.3, fill=PANEL)
    _add_text(s, "LLM providers integrated",
              left=7.05, top=1.7, width=5.6, height=0.4, size=15, bold=True, color=AMBER)
    rows = [
        ("Anthropic Claude (direct)",  "tool_use · strict schemas"),
        ("AWS Bedrock",                "Converse API · same Claude family"),
        ("Azure OpenAI",                "tool_calls · deployments"),
        ("OpenAI",                       "tool_calls · function schemas"),
        ("OpenRouter",                   "runtime schema probe + cache"),
    ]
    y = 2.2
    for k, v in rows:
        _add_text(s, k, left=7.05, top=y, width=3.4, height=0.32, size=12, bold=True, color=TEXT)
        _add_text(s, v, left=10.4, top=y, width=2.4, height=0.32, size=11,
                  color=MUTED, align=PP_ALIGN.RIGHT)
        y += 0.55
    _add_text(s, "Cross-product regression matrix lands by June 13",
              left=7.05, top=5.4, width=5.6, height=0.32, size=11, color=AMBER)

    _footer(s, "Agentic · Conversations & Providers",
            "All 5 providers wired · matrix tests in progress")
    _notes(s, "Provider portability is what makes this safe to deploy in different agencies — "
              "Bedrock for AWS-native, Azure for Microsoft-native, Anthropic direct for "
              "everyone else. The agent code never branches on provider; the translator layer "
              "absorbs the differences.")

    # ── Slide 16: Section divider — Surface & Ops ─────────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _section_divider(s,
        "Section 5 of 5",
        "Surface & Operations",
        "API, frontend, settings, multi-dialect metadata, eval harness — the parts an operator touches",
        color=GREEN)
    _notes(s, "Everything an operator or user actually interacts with day-to-day.")

    # ── Slide 17: API + Frontend ──────────────────────────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_title(s, "API surface + browser experience",
                  "18 endpoints · WebSocket + SSE streaming · 7 browser pages")

    _panel(s, left=0.5, top=1.6, width=6.0, height=5.3, fill=PANEL)
    _add_text(s, "REST + streaming endpoints",
              left=0.7, top=1.7, width=5.6, height=0.4, size=15, bold=True, color=ACCENT)
    _bullets(s, [
        "/health · status, active provider, dialect, catalog readiness",
        "/query, /query/stream — synchronous + WebSocket pipeline",
        "/chat, /chat/stream — agent loop sync + SSE streaming",
        "/conversations — multi-turn lifecycle (CRUD + dialect badges)",
        "/tables, /tables/{fqn}, /domains — catalog browse",
        "/gold — CRUD + approve/reject + retrieval refresh",
        "/admin — config, jobs/rebuild SSE, test_metadata_db",
    ], left=0.7, top=2.15, width=5.6, height=4.5, size=11)

    _panel(s, left=6.85, top=1.6, width=6.0, height=5.3, fill=PANEL)
    _add_text(s, "Browser pages (Next.js)",
              left=7.05, top=1.7, width=5.6, height=0.4, size=15, bold=True, color=AMBER)
    _bullets(s, [
        "Query — single-shot NL → SQL with chart & summary",
        "Chat — multi-turn agent loop with conversation list",
        "Tables — schema browser, FQN detail, sample rows",
        "Domains — domain list with table counts",
        "Gold — gold-query CRUD, approve workflow, eval drilldown",
        "Settings — DB / LLM / embedding connector forms (editable)",
        "Onboarding — sitewide banner when artifacts not yet built",
    ], left=7.05, top=2.15, width=5.6, height=4.5, size=11)

    _footer(s, "Surface · API + Frontend",
            "test_component11 + test_component12 green · provider-aware refresh")
    _notes(s, "Provider-aware refresh (P1/P2) is the secret sauce: switch the active database "
              "in Settings, the Tables and Domains pages refresh automatically — no stale state.")

    # ── Slide 18: Operations — settings, multi-dialect, eval, demo DB ─────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_title(s, "Operator-grade plumbing",
                  "Settings UI · multi-dialect metadata · eval harness · zero-infra demo target")

    panels = [
        ("Settings UI",
         ["DB / LLM / embedding connector forms with inline test buttons",
          "Secrets stored in gitignored runtime_secrets.json (not in body)",
          "Runtime overlay merged on top of default.yaml at boot",
          "Rebuild orchestrator: queue stages, watch live log via SSE"]),
        ("Multi-dialect metadata DB",
         ["SQLite (default zero-infra), Postgres, MSSQL all supported",
          "URL.create() escaping handles passwords with @ : / ?",
          "Forward-only ALTER TABLE migrations for legacy installs",
          "Password redacted from driver error responses"]),
        ("Eval harness",
         ["Gold-question suite with JSON + Markdown reports",
          "Schema-hit, join-hit, descriptor-leakage, BLEU, exact-match metrics",
          "Per-build comparison so regressions surface immediately",
          "CI gate hook in place; expanding to 100 questions for v0.9"]),
        ("Zero-infra demo target",
         ["sample_demo.sqlite shipped with the repo (Ed-Fi-shaped, ~70 rows)",
          "git clone + text2sql serve works without an operator-supplied DB",
          "Reflection picks up unknown tables on operator-supplied DBs",
          "Onboarding banner walks new users to first query"]),
    ]
    rows_per = 2
    cell_w, cell_h = 6.15, 2.6
    for i, (title, items) in enumerate(panels):
        col = i % rows_per
        row = i // rows_per
        x = 0.5 + col * (cell_w + 0.05)
        y = 1.6 + row * (cell_h + 0.1)
        _panel(s, left=x, top=y, width=cell_w, height=cell_h, fill=PANEL)
        _add_text(s, title, left=x + 0.2, top=y + 0.1, width=cell_w - 0.4, height=0.4,
                  size=14, bold=True, color=ACCENT)
        _bullets(s, items, left=x + 0.2, top=y + 0.5, width=cell_w - 0.4,
                 height=cell_h - 0.6, size=10)

    _footer(s, "Operations · Settings + Multi-dialect + Eval + Demo target",
            "All four committed and tested")
    _notes(s, "These are the pieces that make the platform actually usable on day one — and the "
              "reason a fresh git clone can demo without a populated database.")

    # ── Slide 19: What ships next (Jun 13 release) ────────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_title(s, "What ships in the v0.9 release",
                  "6 working weeks of focused polish · then tag and announce")

    items = [
        ("Auto sub-clustering (Leiden)",
         "AMBER",
         "Refines oversize Ed-Fi domains (Student, Assessment, Discipline) into 8–20-table sub-clusters. "
         "Today: static taxonomy works. Next: sharper retrieval inside the heaviest domains."),
        ("Multi-provider regression matrix",
         "AMBER",
         "All 5 LLM providers + 4 embedding providers wired. Matrix expands to test "
         "every consumer (chat, query, eval, agent) against every provider on every supported dialect."),
        ("Agentic polish",
         "AMBER",
         "Tool-call reliability under streaming, richer agent tools (chart export, eval drilldown), "
         "smoother conversation list UX, retry on transient provider errors."),
        ("Observability (OTel + Prometheus)",
         "RED",
         "Per-stage tracing spans + counters + Grafana dashboards. Logging is in place today; "
         "spans and metrics are next. May slip to v0.10 if matrix work runs long."),
        ("Recorded demo + release cut",
         "AMBER",
         "End-to-end recording on a clean Mac and Windows machine; v0.9 tag pushed Fri June 13; "
         "release notes + operator runbook + provider-swap guide finalized."),
    ]
    _panel(s, left=0.5, top=1.6, width=12.3, height=5.4, fill=PANEL)
    y = 1.75
    for label, level, what in items:
        color = AMBER if level == "AMBER" else RED
        _add_text(s, "●", left=0.7, top=y, width=0.3, height=0.3, size=18, bold=True, color=color)
        _add_text(s, label, left=1.05, top=y, width=11.7, height=0.3, size=13, bold=True, color=TEXT)
        _add_text(s, what, left=1.05, top=y + 0.32, width=11.7, height=0.7, size=11, color=MUTED)
        y += 1.05

    _footer(s, "Plan · Section 6", "")
    _notes(s, "The honest framing: most of the platform is already built. The June 13 date buys "
              "real headroom to land the things that matter for a confident demo — provider "
              "matrix, agent reliability, observability, performance, and a recorded walkthrough. "
              "If observability runs long it slips to v0.10 cleanly without affecting the demo.")

    # ── Slide 20: Phased plan to release ──────────────────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_title(s, "Phased plan to v0.9",
                  "Six working weeks · Apr 30 → Fri June 12 · release tag Fri June 13")

    phases = [
        ("Week 1",  "Apr 30 + May 4–8",
         "Provider matrix + Leiden — kickoff",
         GREEN,
         ["Lock matrix shape: 5 LLMs × 4 embeddings × 3 dialects × 4 consumers; CI fixtures",
          "Run agent + chat against Anthropic, Bedrock first; patch translator gaps",
          "Spike Leiden auto-clustering on oversize domains (Student / Assessment / Discipline)",
          "Cluster-ID stability via Hungarian assignment across rebuilds"]),
        ("Week 2",  "May 11–15",
         "Provider matrix coverage + agentic polish",
         GREEN,
         ["Add Azure OpenAI, OpenAI, OpenRouter to the matrix; capture per-provider regressions",
          "Tool-call retries + structured failure surfaces in the agent loop",
          "New agent tools: chart export, eval drilldown, conversation rename / archive",
          "Streaming UX polish — token-level feedback under flaky network conditions"]),
        ("Week 3",  "May 18–22",
         "Observability scaffolding",
         AMBER,
         ["OpenTelemetry spans on every pipeline + agent stage; correlation IDs across services",
          "Prometheus counters / histograms; Grafana dashboards in infra/grafana/",
          "Structured logging context: request_id, provider, dialect, conversation_id",
          "Stretch — slips to v0.10 without affecting the demo if the matrix runs long"]),
        ("Week 4",  "May 25–29",
         "Performance pass + eval suite expansion",
         GREEN,
         ["Cache hot paths: APSP load, retrieval, gold-store ANN; meet p50/p95 budget",
          "Per-user quota + rate limits; secret-leak scrubber on every error path",
          "Expand gold-query suite from 50 → 100 questions; wire nightly CI gate",
          "Eval Dashboard polish — per-build deltas, regression alerts, drilldown"]),
        ("Week 5",  "Jun 1–5",
         "Cross-platform smoke + documentation",
         GREEN,
         ["Clean Windows 11 + macOS dry runs from git clone; capture install-time issues",
          "Provider-swap drill: Anthropic → Bedrock → Azure mid-session without restart",
          "Operator runbook, provider-swap guide, gold-curation guide — all finalized",
          "Release notes draft against the v0.9 candidate"]),
        ("Week 6",  "Jun 8–12",
         "Recorded demo + release prep",
         GREEN,
         ["Demo script: connect demo SQLite → ask question → chart + summary → switch DB",
          "Re-record any segments where the live system stutters; caption + edit pass",
          "Final eval gate against v0.9 RC; capture metrics for release notes",
          "v0.9 release-candidate tagged + smoke-tested on clean machines"]),
        ("Fri Jun 13",  "Release day",
         "v0.9 RELEASE",
         ACCENT,
         ["Tag v0.9 on main; push to GitHub Releases with the recorded demo attached",
          "Distribute client-facing release summary + recorded walkthrough",
          "Internal announcement; v0.10 backlog opened (observability stretch + extras)"]),
    ]

    _panel(s, left=0.5, top=1.55, width=12.3, height=5.5, fill=PANEL)
    y = 1.65
    for week_label, dates, theme, color, tasks in phases:
        is_release = week_label == "Fri Jun 13"
        _add_text(s, week_label, left=0.7, top=y, width=1.2, height=0.26,
                  size=11, bold=True, color=color)
        _add_text(s, dates, left=1.95, top=y, width=1.6, height=0.26,
                  size=10, color=MUTED)
        _add_text(s, theme, left=3.6, top=y, width=9.2, height=0.26,
                  size=11, bold=True, color=TEXT if not is_release else GREEN)
        for t in tasks:
            y += 0.19
            _add_text(s, "•  " + t, left=3.6, top=y, width=9.2, height=0.22,
                      size=9, color=MUTED)
        y += 0.16

    _footer(s, "Plan · 6 working weeks + release day",
            "May 1 holiday + weekends excluded")
    _notes(s, "Why six weeks (not eleven days): the punch list is real engineering work, not "
              "checklist items. Provider matrix touches 60 distinct combinations and historically "
              "surfaces translator-layer bugs that take a half-day each. Leiden tuning needs "
              "iteration. Observability work touches every span in the system. Recording a "
              "client-facing demo we'd actually ship takes time and re-takes. Six weeks gives "
              "honest buffer; if the team finishes early we ship early — but we won't book a "
              "heroics-only plan.")

    # ── Slide 21: Risks ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_title(s, "Risks & how we'll handle them",
                  "Known unknowns flagged early; fallbacks scoped into the plan")

    risks = [
        ("Provider matrix surprises",
         "A new provider/consumer combination may need translator-layer fixes that take longer than a half-day.",
         "Two-day budget covers worst case; OpenRouter's runtime probe absorbs edge cases automatically.",
         AMBER),
        ("Leiden tuning quality",
         "Sub-cluster quality varies with affinity-matrix weights; may not converge cleanly on Day 3.",
         "Static taxonomy stays as the v0.9 default; Leiden ships behind a feature flag if needed.",
         AMBER),
        ("Observability slippage",
         "OTel + Prometheus spans across 9 pipeline stages + 8 agent tools is non-trivial.",
         "Stretch goal — slips to v0.10 cleanly. Logging is comprehensive today, so debug-ability holds.",
         RED),
        ("Live demo connectivity",
         "Live demos can stutter on conference Wi-Fi or provider rate limits.",
         "Pre-record the full flow in Week 6 + polish before release-day; live Q&A only.",
         GREEN),
    ]
    _panel(s, left=0.5, top=1.6, width=12.3, height=5.4, fill=PANEL)
    y = 1.8
    for label, what, plan, color in risks:
        _add_text(s, "●", left=0.7, top=y, width=0.3, height=0.3, size=20, bold=True, color=color)
        _add_text(s, label, left=1.1, top=y, width=11.7, height=0.3, size=14, bold=True, color=TEXT)
        _add_text(s, "Risk: "  + what, left=1.1, top=y + 0.32, width=11.7, height=0.32, size=11, color=MUTED)
        _add_text(s, "Plan: "  + plan, left=1.1, top=y + 0.62, width=11.7, height=0.32, size=11, color=ACCENT)
        y += 1.27

    _footer(s, "Risk register", "")
    _notes(s, "We're not promising perfection by Friday June 13. We are promising a tested "
              "platform with a recorded demo, the multi-provider matrix passing, and stretch "
              "items called out before they slip silently. The 6-week timeline is honest "
              "engineering scope, not heroics.")

    # ── Slide 22: What you'll see ─────────────────────────────────────────
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_text(s, "What you'll see at the v0.9 demo (June 13)", left=0.8, top=1.4, width=11.7, height=0.6,
              size=24, bold=True, color=ACCENT)

    _panel(s, left=0.5, top=2.1, width=12.3, height=4.6, fill=PANEL)
    _bullets(s, [
        "Connect any Ed-Fi-shaped database from a browser — no code, no config files",
        "Ask plain-English questions and watch SQL + rows + chart + summary land in seconds",
        "Hold a multi-turn conversation — follow-up questions reuse earlier context",
        "Watch the agent stream tool calls live: classify → retrieve → resolve → run",
        "See the SQL color-coded against the catalog with cited tables",
        "Approve a query — and watch the platform get smarter for the next question like it",
        "Switch from one database (or LLM provider) to another mid-session and see results refresh automatically",
        "Open the Eval Dashboard to see how accuracy is trending across builds",
    ], left=0.7, top=2.25, width=11.9, height=4.4, size=13, color=TEXT)

    _add_text(s, "Demo machine: clean Windows install · no prior setup · v0.9 tagged Fri June 13",
              left=0.5, top=6.85, width=12.3, height=0.4, size=12, color=MUTED, align=PP_ALIGN.CENTER)

    _notes(s, "End on the demo. The 8 specific things on this slide are what you'll watch. "
              "Each bullet is gated on a workstream above. We sequence them so even if observability "
              "slips, the first 7 are still demo-ready.")

    # ── Save ──────────────────────────────────────────────────────────────
    out = Path(__file__).resolve().parent / "architecture_deck.pptx"
    prs.save(str(out))
    print(f"Wrote {out}  ({out.stat().st_size / 1024:.1f} KB · {len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
